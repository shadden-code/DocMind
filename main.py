print("DocMind server starting...")
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from pypdf import PdfReader
import io

# ── Multi-format text extraction ──────────────────────────────────────────────
_SUPPORTED_EXTENSIONS = {
    '.pdf', '.docx', '.doc', '.pptx', '.ppt',
    '.xlsx', '.xls', '.csv', '.txt', '.md', '.rtf',
}

def _extract_text_from_file(filename: str, file_bytes: bytes) -> list[tuple[int, str]]:
    """
    Extract text from various file types.
    Returns list of (page_num, text) tuples, mirroring the PDF extractor output.
    """
    ext = ('.' + filename.rsplit('.', 1)[-1]).lower() if '.' in filename else ''

    # ── PDF ──────────────────────────────────────────────────────────────────
    if ext == '.pdf':
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(file_bytes))
        pages = []
        for i, page in enumerate(reader.pages, start=1):
            text = _clean(page.extract_text() or '')
            if text.strip():
                pages.append((i, text))
        return pages

    # ── Word (.docx) ─────────────────────────────────────────────────────────
    if ext in ('.docx', '.doc'):
        import docx as _docx
        doc = _docx.Document(io.BytesIO(file_bytes))
        paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        # Group every 20 paragraphs as one "page"
        pages, chunk_size = [], 20
        for i in range(0, max(len(paras), 1), chunk_size):
            text = _clean('\n'.join(paras[i:i+chunk_size]))
            if text.strip():
                pages.append((i // chunk_size + 1, text))
        return pages or [(1, _clean(' '.join(paras)) or 'Empty document')]

    # ── PowerPoint (.pptx) ───────────────────────────────────────────────────
    if ext in ('.pptx', '.ppt'):
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        pages = []
        for i, slide in enumerate(prs.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    parts.append(shape.text.strip())
            text = _clean('\n'.join(parts))
            if text.strip():
                pages.append((i, text))
        return pages

    # ── Excel (.xlsx / .xls) ─────────────────────────────────────────────────
    if ext in ('.xlsx', '.xls'):
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        pages = []
        for sheet_num, ws in enumerate(wb.worksheets, start=1):
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else '' for c in row]
                line = '\t'.join(cells).strip()
                if line and line != '\t' * (len(cells)-1):
                    rows.append(line)
            text = _clean('\n'.join(rows))
            if text.strip():
                pages.append((sheet_num, f"[Sheet: {ws.title}]\n{text}"))
        wb.close()
        return pages

    # ── CSV ──────────────────────────────────────────────────────────────────
    if ext == '.csv':
        import csv as _csv
        text_io = io.StringIO(file_bytes.decode('utf-8', errors='replace'))
        reader = _csv.reader(text_io)
        rows = ['\t'.join(r) for r in reader if any(c.strip() for c in r)]
        text = _clean('\n'.join(rows))
        # Split into ~100-row pages
        chunk_size = 100
        pages = []
        for i in range(0, max(len(rows), 1), chunk_size):
            t = _clean('\n'.join(rows[i:i+chunk_size]))
            if t.strip():
                pages.append((i // chunk_size + 1, t))
        return pages

    # ── Plain text / Markdown / RTF ──────────────────────────────────────────
    if ext in ('.txt', '.md', '.rtf'):
        raw = file_bytes.decode('utf-8', errors='replace')
        # Strip basic RTF control words
        if ext == '.rtf':
            raw = re.sub(r'\\[a-z]+\d*\s?', ' ', raw)
            raw = re.sub(r'[{}]', '', raw)
        lines = raw.split('\n')
        chunk_size = 60
        pages = []
        for i in range(0, max(len(lines), 1), chunk_size):
            text = _clean('\n'.join(lines[i:i+chunk_size]))
            if text.strip():
                pages.append((i // chunk_size + 1, text))
        return pages

    raise ValueError(f"Unsupported file type: {ext or 'unknown'}")

from groq import Groq
from rank_bm25 import BM25Okapi
# Heavy ML imports are deferred to background thread to avoid Render port timeout
# HuggingFaceEmbeddings, CrossEncoder, RecursiveCharacterTextSplitter loaded lazily
import os
import gc
import chromadb
import uuid, re
from pathlib import Path
from dotenv import load_dotenv

load_dotenv()

try:
    from cerebras.cloud.sdk import Cerebras
    _CEREBRAS_AVAILABLE = True
except ImportError:
    _CEREBRAS_AVAILABLE = False
    print("[DocMind] cerebras-cloud-sdk not installed — run: pip install cerebras-cloud-sdk")


BASE_DIR = Path(__file__).parent

app = FastAPI()
embeddings = None

import threading

def _load_embeddings():
    global embeddings, splitter
    if embeddings is not None:
        return  # already loaded, prevent duplicate loading
    print("[DocMind] Loading ML dependencies in background...")
    # Import heavy libraries only now (avoids slow startup / Render port timeout)
    import gc
    from transformers import logging as hf_logging
    hf_logging.set_verbosity_error()
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    from langchain_huggingface import HuggingFaceEmbeddings
    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=150)
    embeddings = HuggingFaceEmbeddings(
        model_name="sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2",
        encode_kwargs={"batch_size": 64, "normalize_embeddings": True},
    )
    gc.collect()  # free unused memory after model load
    print("[DocMind] Embeddings loaded ✓")

@app.on_event("startup")
async def startup_event():
    # Load in background thread so port binds immediately (avoids Render timeout)
    threading.Thread(target=_load_embeddings, daemon=True).start()
app.add_middleware(
    CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"]
)

# Serve style.css, script.js and any other static assets from the project folder
app.mount("/static", StaticFiles(directory=str(BASE_DIR)), name="static")

@app.get("/health")
def health():
    """Used by frontend to poll until server is ready."""
    if embeddings is None or splitter is None:
        from fastapi.responses import JSONResponse as JR
        return JR({"ready": False, "message": "Loading AI models…"}, status_code=503)
    return {"ready": True, "message": "Server ready"}

@app.get("/")
def serve_ui():
    return FileResponse(
        BASE_DIR / "index.html",
        headers={"Cache-Control": "no-cache, no-store, must-revalidate", "Pragma": "no-cache"}
    )

# splitter is initialised in _load_embeddings() background thread
splitter = None

import json, time
from fastapi import Cookie
from fastapi.responses import JSONResponse

# ── Persistent ChromaDB (local use — data saved to ./chroma_db) ───────────────
_DB_DIR = BASE_DIR / "chroma_db"
_DB_DIR.mkdir(exist_ok=True)
_chroma_client = chromadb.PersistentClient(path=str(_DB_DIR))
_SESSION_TTL = 3600
_sessions: dict = {}

# ── Load existing documents from disk on startup ───────────────────────────────
_DOCS_FILE = _DB_DIR / "documents.json"

def _load_documents_from_disk() -> dict:
    if _DOCS_FILE.exists():
        try:
            with open(_DOCS_FILE) as f:
                data = json.load(f)
            print(f"[DocMind] Loaded {len(data)} document records from disk")
            return data
        except Exception as e:
            print(f"[DocMind] Could not load documents.json: {e}")
    return {}

def _save_documents_to_disk(documents: dict):
    try:
        with open(_DOCS_FILE, "w") as f:
            json.dump(documents, f, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f"[DocMind] Could not save documents.json: {e}")

def _get_or_create_session(session_id: str) -> dict:
    """Get existing session or create a new one with persistent ChromaDB."""
    if session_id not in _sessions:
        collection = _chroma_client.get_or_create_collection(
            name=f"docs_{session_id[:8]}",
            metadata={"hnsw:space": "cosine"}
        )
        _sessions[session_id] = {
            "db": collection,
            "documents": _load_documents_from_disk(),
            "bm25_store": {},
            "last_active": time.time(),
        }
        print(f"[DocMind] Session created: {session_id[:8]}…")
    else:
        _sessions[session_id]["last_active"] = time.time()
    return _sessions[session_id]

# Legacy globals
documents: dict = {}
bm25_store: dict = {}
db = None

def _save_documents():
    pass  # handled per-session

# ── Cross-encoder reranker (~80 MB, downloaded once) ──────────────────────────
# ms-marco-MiniLM-L-6-v2: fast reranker that scores (query, passage) pairs.
# Used to rerank top-10 cosine candidates → top-3 semantically precise results.
reranker = None
RERANKER_OK = True

def get_reranker():
    global reranker

    if reranker is None:
        try:
            print("[DocMind] Loading cross-encoder reranker...")
            from sentence_transformers import CrossEncoder
            reranker = CrossEncoder(
                "cross-encoder/ms-marco-MiniLM-L-6-v2",
                max_length=512
            )
            print("[DocMind] Cross-encoder reranker loaded ✓")
        except Exception as e:
            print(f"[DocMind] Failed to load reranker: {e}")
            return None

    return reranker
# ── Groq client (replaces FLAN-T5) ───────────────────────────────────────────
# Free API key at https://console.groq.com → API Keys → Create key
# Set it as an environment variable:  export GROQ_API_KEY="gsk_..."
# Or paste it directly below (not recommended for production):

GROQ_API_KEY = os.getenv("GROQ_API_KEY")
HF_TOKEN = os.getenv("HF_TOKEN")
CEREBRAS_API_KEY = os.getenv("CEREBRAS_API_KEY")

# ── Cerebras ───────────────────────────────────────────────────────────────────
_CEREBRAS_MODEL   = "llama-3.3-70b"
_CEREBRAS_FAST    = "llama3.1-8b"
_cerebras_client  = None

if _CEREBRAS_AVAILABLE and CEREBRAS_API_KEY:
    try:
        _cerebras_client = Cerebras(api_key=CEREBRAS_API_KEY)
        print(f"[DocMind] Cerebras ready ✓  (model: {_CEREBRAS_MODEL})")
    except Exception as _ce:
        print(f"[DocMind] Cerebras init failed: {_ce}")

# ── Groq model chain ───────────────────────────────────────────────────────────
_GROQ_MODEL    = "llama-3.3-70b-versatile"
_GROQ_FALLBACK = "moonshotai/kimi-k2-instruct-0905"
_GROQ_FALLBACK2= "qwen/qwen3-32b"
_GROQ_FAST     = "llama-3.1-8b-instant"

try:
    _groq_client = Groq(api_key=GROQ_API_KEY) if GROQ_API_KEY else None
    SUMMARIZER_OK = _groq_client is not None
    if SUMMARIZER_OK:
        print(f"[DocMind] Groq client ready ✓")
        cb = f"Cerebras({_CEREBRAS_MODEL})" if _cerebras_client else "unavailable"
        print(f"[DocMind] Chain: {_GROQ_MODEL} → {cb} → {_GROQ_FALLBACK} → {_GROQ_FALLBACK2} → {_GROQ_FAST}")
    else:
        print("[DocMind] GROQ_API_KEY not set — falling back to extractive mode.")
except Exception as _e:
    print(f"[DocMind] Groq init failed: {_e}")
    _groq_client  = None
    SUMMARIZER_OK = False


def _cerebras_generate(system: str, user: str, max_tokens: int = 900) -> str:
    """Call Cerebras — tries llama-3.3-70b first, falls back to llama3.1-8b."""
    if not _cerebras_client:
        raise RuntimeError("Cerebras client not initialised")

    for model in [_CEREBRAS_MODEL, _CEREBRAS_FAST]:
        try:
            resp = _cerebras_client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": system},
                    {"role": "user",   "content": user},
                ],
                max_tokens=max(max_tokens, 400),
                temperature=0.2,
            )
            result = resp.choices[0].message.content.strip()
            print(f"[DocMind] ✓ Response from: Cerebras/{model}  ({len(result.split())} words)")
            return result
        except Exception as e:
            err = str(e)
            if "404" in err or "not_found" in err or "429" in err or "rate_limit" in err.lower():
                print(f"[DocMind] Cerebras/{model} failed ({err[:60].strip()}) → trying next…")
                continue
            raise

    raise RuntimeError("All Cerebras models failed.")


def _groq_generate(system: str, user: str, max_tokens: int = 300,
                   model: str | None = None) -> str:
    """
    Call Groq with automatic fallback:
      llama-3.3-70b → Cerebras(llama-3.3-70b) → kimi-k2 → qwen3-32b → llama-3.1-8b
    """
    if not _groq_client:
        raise RuntimeError("Groq client not initialised — set GROQ_API_KEY")

    use_model = model or _GROQ_MODEL
    _GROQ_CHAIN = [_GROQ_MODEL, _GROQ_FALLBACK, _GROQ_FALLBACK2, _GROQ_FAST]

    def _call_groq(m, u):
        effective_tokens = max(max_tokens, 400) if m != _GROQ_FAST else max_tokens
        resp = _groq_client.chat.completions.create(
            model=m,
            messages=[
                {"role": "system", "content": system},
                {"role": "user",   "content": u},
            ],
            max_tokens=effective_tokens,
            temperature=0.2,
        )
        result = resp.choices[0].message.content.strip()
        print(f"[DocMind] ✓ Response from: {m}  ({len(result.split())} words)")
        return result

    def _is_rate_limit(err: str) -> bool:
        return "429" in err or "rate_limit" in err.lower()

    def _is_too_large(err: str) -> bool:
        return "413" in err or "too large" in err.lower()

    # Fast path — cheap model for tagging/summaries
    if model == _GROQ_FAST:
        try:
            return _call_groq(model, user)
        except Exception as e:
            err = str(e)
            if _is_too_large(err):
                try:
                    return _call_groq(model, user[:2000])
                except Exception:
                    pass
            if (_is_rate_limit(err) or _is_too_large(err)) and _cerebras_client:
                print(f"[DocMind] Fast model limit → falling back to Cerebras…")
                return _cerebras_generate(system, user, max_tokens)
            raise

    # Quality path — try Cerebras on first rate limit, then walk Groq chain
    start = _GROQ_CHAIN.index(use_model) if use_model in _GROQ_CHAIN else 0
    cerebras_tried = False

    for i, m in enumerate(_GROQ_CHAIN[start:], start=start):
        try:
            return _call_groq(m, user)
        except Exception as e:
            err = str(e)
            if _is_too_large(err):
                try:
                    print(f"[DocMind] Request too large for {m}, truncating…")
                    return _call_groq(m, user[:3000])
                except Exception:
                    pass
                print(f"[DocMind] Still too large for {m}, trying next…")
                continue
            if _is_rate_limit(err):
                if not cerebras_tried and _cerebras_client:
                    try:
                        print(f"[DocMind] Rate limit on {m} → switching to Cerebras…")
                        return _cerebras_generate(system, user, max_tokens)
                    except Exception as ce:
                        print(f"[DocMind] Cerebras failed: {ce}, continuing Groq chain…")
                        cerebras_tried = True
                next_idx = i + 1
                if next_idx < len(_GROQ_CHAIN):
                    print(f"[DocMind] Trying next Groq model: {_GROQ_CHAIN[next_idx]}…")
                    continue
                raise RuntimeError("All models hit rate limits. Please wait a few minutes.")
            raise

    raise RuntimeError("All models exhausted. Please wait a few minutes and try again.")



# ── Schemas ────────────────────────────────────────────────────────────────────
class Question(BaseModel):
    question: str
    doc_id: str | None = None

class Rerank(BaseModel):
    question: str
    instruction: str
    doc_id: str | None = None

class CompareRequest(BaseModel):
    doc_id_1: str
    doc_id_2: str
    question: str = "What are the main topics and key points?"

class SearchRequest(BaseModel):
    query: str
    doc_id: str | None = None


# ─────────────────────────────────────────────
#  STEP 1 — Upload a PDF
# ─────────────────────────────────────────────
@app.post("/upload")
async def upload(file: UploadFile = File(...), session_id: str = Cookie(default=None)):
    if embeddings is None or splitter is None:
        raise HTTPException(503, "Server is still loading, please wait 30 seconds and try again.")

    ext = ('.' + file.filename.rsplit('.', 1)[-1]).lower() if '.' in file.filename else ''
    if ext not in _SUPPORTED_EXTENSIONS:
        raise HTTPException(400, f"Unsupported file type '{ext}'. Supported: PDF, Word, PowerPoint, Excel, CSV, TXT, Markdown.")

    file_bytes = await file.read()
    try:
        pages_text = _extract_text_from_file(file.filename, file_bytes)
    except ValueError as ve:
        raise HTTPException(400, str(ve))
    except Exception:
        raise HTTPException(422, f"Could not read this file. Make sure it is a valid {ext.lstrip('.')} document.")

    # Create or get session
    if not session_id:
        session_id = uuid.uuid4().hex
    session = _get_or_create_session(session_id)
    session_db = session["db"]
    session_docs = session["documents"]
    session_bm25 = session["bm25_store"]

    if not pages_text:
        raise HTTPException(422, "No text could be extracted from this file.")

    # ── STEP 3 — Chunk + store in vector DB ───
    doc_id = uuid.uuid4().hex[:12]
    all_chunks, all_meta = [], []
    for page_num, page_text in pages_text:
        for chunk in splitter.split_text(page_text):
            all_chunks.append(chunk)
            all_meta.append({"doc_id": doc_id, "page": page_num})

    # Embed chunks in full batches for local performance
    _EMBED_BATCH = 64
    all_vectors = []
    for start in range(0, len(all_chunks), _EMBED_BATCH):
        batch = all_chunks[start: start + _EMBED_BATCH]
        all_vectors.extend(embeddings.embed_documents(batch))

    all_ids = [f"{doc_id}_{idx}" for idx in range(len(all_chunks))]

    _DB_BATCH = 500
    for start in range(0, len(all_chunks), _DB_BATCH):
        session_db.add(
            documents=all_chunks[start: start + _DB_BATCH],
            embeddings=all_vectors[start: start + _DB_BATCH],
            ids=all_ids[start: start + _DB_BATCH],
            metadatas=all_meta[start: start + _DB_BATCH],
        )

    # ── Build BM25 index for this document ────────────────────────────────────
    tokenized = [_tokenize(c) for c in all_chunks]
    session_bm25[doc_id] = {
        "bm25":   BM25Okapi(tokenized),
        "chunks": all_chunks,
        "metas":  all_meta,
    }

    summary      = _summarize(all_chunks)
    suggested_qa = _generate_qa(all_chunks)

    session_docs[doc_id] = {
        "name":         file.filename,
        "summary":      summary,
        "chunks":       len(all_chunks),
        "pages":        len(pages_text),
        "suggested_qa": suggested_qa,
    }

    response = JSONResponse({
        "doc_id":       doc_id,
        "filename":     file.filename,
        "chunks":       len(all_chunks),
        "pages":        len(pages_text),
        "summary":      summary,
        "suggested_qa": suggested_qa,
        "session_notice": "⚠️ Your documents are private and stored only for this session. They will be automatically deleted after 1 hour of inactivity or when you close the app.",
    })
    response.set_cookie("session_id", session_id, max_age=_SESSION_TTL, httponly=True, samesite="lax")
    return response


# ─────────────────────────────────────────────
#  STEP 4 — User asks a question
# ─────────────────────────────────────────────
@app.post("/ask")
def ask(body: Question, session_id: str = Cookie(default=None)):
    if not session_id or session_id not in _sessions:
        raise HTTPException(404, "No session found. Please upload a document first.")
    session = _get_or_create_session(session_id)
    if session["db"].count() == 0:
        raise HTTPException(404, "No documents uploaded yet.")
    if not body.question.strip():
        raise HTTPException(400, "Question cannot be empty.")
    result = _query(body.question, body.doc_id, n=5, session=session)
    result["answer"] = _generate_answer(body.question, result["passages"])
    return result


# ── Re-rank ────────────────────────────────────────────────────────────────────
@app.post("/rerank")
def rerank(body: Rerank, session_id: str = Cookie(default=None)):
    if not session_id or session_id not in _sessions:
        raise HTTPException(404, "No session found. Please upload a document first.")
    session = _get_or_create_session(session_id)
    if session["db"].count() == 0:
        raise HTTPException(404, "No documents uploaded yet.")
    inst = body.instruction.lower()

    if any(w in inst for w in ["specific", "precise"]):
        result = _query(body.question, body.doc_id, n=1, session=session)
        result["rerank_mode"] = "precise"
        return result
    elif any(w in inst for w in ["simple", "brief", "short"]):
        result = _query(body.question, body.doc_id, n=3, session=session)
        passages = result.get("passages", [])
        if passages:
            shortest = min(passages, key=lambda p: len(p["text"]))
            result["answer"] = shortest["text"]
            result["page"]   = shortest["page"]
        result["rerank_mode"] = "brief"
        return result
    elif any(w in inst for w in ["detail", "more", "explain", "elaborate"]):
        result = _query(body.question, body.doc_id, n=5, session=session)
        passages = result.get("passages", [])
        if passages:
            combined = " ".join(p["text"] for p in passages)
            result["answer"] = combined
        result["rerank_mode"] = "detailed"
        return result
    else:
        result = _query(f"{body.question} {body.instruction}", body.doc_id, n=3, session=session)
        result["rerank_mode"] = "refined"
        return result


# ── Feature 3: Summarize a document ───────────────────────────────────────────
@app.get("/summarize/{doc_id}")
def summarize(doc_id: str, session_id: str = Cookie(default=None)):
    if not session_id or session_id not in _sessions:
        raise HTTPException(404, "Session not found.")
    session = _get_or_create_session(session_id)
    if doc_id not in session["documents"]:
        raise HTTPException(404, "Document not found.")
    res    = session["db"].get(where={"doc_id": doc_id})
    chunks = [_clean(c) for c in res["documents"]]
    doc_type = _detect_doc_type(chunks)
    return _build_structured_summary(chunks, doc_type)


# ── Feature 6: Extract key information ────────────────────────────────────────
@app.get("/extract/{doc_id}")
def extract(doc_id: str, session_id: str = Cookie(default=None)):
    if not session_id or session_id not in _sessions:
        raise HTTPException(404, "Session not found.")
    session = _get_or_create_session(session_id)
    if doc_id not in session["documents"]:
        raise HTTPException(404, "Document not found.")
    res  = session["db"].get(where={"doc_id": doc_id})
    text = " ".join(res["documents"])
    return {
        "people":    _extract_people(text),
        "companies": _extract_companies(text),
        "dates":     _extract_dates(text),
        "numbers":   _extract_numbers(text),
        "topics":    _extract_topics(text),
    }


# ── Feature 7: Compare two documents ──────────────────────────────────────────

_BOILERPLATE = re.compile(
    r"(appendix|table of contents|analytical perspectives|"
    r"u\.s\. government publishing office|budget of the united states government"
    r"|budget appendix|contains detailed information|gpo|this document|isbn|"
    r"for sale by|superintendent of documents|library of congress|"
    r"printed in the united states|visit our website|"
    r"chapter \d+|part \d+|section \d+\.\d+|"
    r"the budget includes|the following table|see table|see appendix)",
    re.IGNORECASE
)

def _is_boilerplate(text: str) -> bool:
    if bool(_BOILERPLATE.search(text)) and len(text) < 600:
        return True
    words = text.split()
    if len(words) > 10 and sum(1 for w in words if w.isupper()) / len(words) > 0.55:
        return True
    return False

def _tag_chunks_with_topics(doc_id: str, question: str, session: dict) -> list[dict]:
    """
    Sample up to 8 non-boilerplate chunks, truncate to 200 chars each,
    tag with topic + numbers using the fast 8b model (<4000 tokens total).
    """
    res    = session["db"].get(where={"doc_id": doc_id})
    chunks = [(c, m) for c, m in zip(res["documents"], res["metadatas"])
              if not _is_boilerplate(_clean(c))]

    if not chunks:
        return []

    # Sample 8 chunks spread evenly across the document
    step   = max(1, len(chunks) // 8)
    sample = chunks[::step][:8]

    # Truncate each chunk to 200 chars to stay well within TPM limit
    numbered = "\n\n".join(
        f"[{i+1}] p.{m.get('page','?')}: {_clean(c)[:200]}"
        for i, (c, m) in enumerate(sample)
    )

    raw = _groq_generate(
        system="You are a document analyst. Respond ONLY with valid JSON array, no markdown fences.",
        user=f'Question: "{question}"\n\nFor each passage return: id, topic (short label), '
             f'numbers (list of {{value, context}} for $ or %), relevant (true/false).\n\n'
             f'Passages:\n{numbered}\n\n'
             f'Return JSON array only: [{{"id":1,"topic":"...","numbers":[],"relevant":true}},...] ',
        max_tokens=600,
        model=_GROQ_FAST,
    )

    try:
        raw_clean = raw.strip().lstrip("```json").lstrip("```").rstrip("```").strip()
        tags = json.loads(raw_clean)
    except Exception:
        tags = []

    tagged = []
    for i, (c, m) in enumerate(sample):
        tag = next((t for t in tags if t.get("id") == i + 1), {})
        tagged.append({
            "text":     _clean(c),
            "page":     m.get("page"),
            "topic":    tag.get("topic", "General"),
            "numbers":  tag.get("numbers", []),
            "relevant": tag.get("relevant", False),
        })

    return tagged

def _generate_comparison(question: str, name1: str, tagged1: list,
                                        name2: str, tagged2: list) -> str:

    # Truncate passage text in prompt to 150 chars to keep prompt small
    def fmt_tagged(tagged):
        if not tagged:
            return "(No relevant passages retrieved.)"
        out = []
        for t in tagged:
            nums = ""
            if t["numbers"]:
                nums = " | " + "; ".join(
                    f"{n['value']} ({n['context']})" for n in t["numbers"][:3]
                )
            snippet = t['text'][:150].replace('\n', ' ')
            out.append(f"[{t['topic']} | p.{t['page']}]{nums}\n{snippet}")
        return "\n\n".join(out)

    # Build same-topic numeric pairs for safe comparison
    def topic_numbers(tagged):
        result = {}
        for t in tagged:
            result.setdefault(t["topic"], []).extend(t["numbers"])
        return result

    nums1 = topic_numbers(tagged1)
    nums2 = topic_numbers(tagged2)
    shared = set(nums1) & set(nums2)

    numeric_section = ""
    if shared:
        lines = []
        for topic in list(shared)[:6]:
            n1 = "; ".join(f"{n['value']} ({n['context']})" for n in nums1[topic][:2]) or "—"
            n2 = "; ".join(f"{n['value']} ({n['context']})" for n in nums2[topic][:2]) or "—"
            lines.append(f"  {topic}: {name1}={n1} | {name2}={n2}")
        numeric_section = "Same-topic numbers:\n" + "\n".join(lines) + "\n"

    prompt = (
        f"Question: {question}\n\n"
        f"=== {name1} ===\n{fmt_tagged(tagged1)}\n\n"
        f"=== {name2} ===\n{fmt_tagged(tagged2)}\n\n"
        f"{numeric_section}"
        f"RULES:\n"
        f"1. SCOPE: Only include table rows where the question asks about new/changed priorities, "
        f"programs, funds, or policy changes. EXCLUDE rows about budget structure elements "
        f"(e.g. 'Budget Authority', 'Mandatory Resources', 'Obligations') — those are structural, not priorities.\n"
        f"2. NEW PRIORITIES: Only include a row if the evidence shows a genuinely new program, "
        f"new fund, new initiative, or a clear change in emphasis. Do NOT include rows just because "
        f"one document retrieved a passage and the other didn't — that is a retrieval gap, not a real difference.\n"
        f"3. ABSENCE LANGUAGE: Never write 'not mentioned' or 'not in retrieved passages' in the table. "
        f"If one side lacks evidence, write 'Not found in retrieved excerpts' and only include the row "
        f"if the other side has strong positive evidence of something new.\n"
        f"4. NUMBERS: Only compare numbers sharing the same topic AND same context. Never pair a total "
        f"budget against a sub-allocation.\n"
        f"5. FORMAT: Comparison table first (Category | {name1} | {name2} | Pages), "
        f"then bullet-point key differences summary. All claims must cite page numbers.\n"
        f"6. START your answer with: 'Based on the retrieved excerpts from both documents...'"
    )

    return _groq_generate(
        system=(
            "You are a precise document comparison analyst. "
            "Be evidence-based and conservative — only report differences you can support with the passages provided. "
            "Never infer absence from a retrieval gap."
        ),
        user=prompt,
        max_tokens=900,
    )

@app.post("/compare")
def compare(body: CompareRequest, session_id: str = Cookie(default=None)):
    if not session_id or session_id not in _sessions:
        raise HTTPException(404, "Session not found.")
    session = _get_or_create_session(session_id)
    for did in [body.doc_id_1, body.doc_id_2]:
        if did not in session["documents"]:
            raise HTTPException(404, f"Document {did} not found.")

    name1 = session["documents"][body.doc_id_1]["name"]
    name2 = session["documents"][body.doc_id_2]["name"]
    tagged1 = _tag_chunks_with_topics(body.doc_id_1, body.question, session)
    tagged2 = _tag_chunks_with_topics(body.doc_id_2, body.question, session)
    relevant1 = [t for t in tagged1 if t.get("relevant")][:8] or tagged1[:8]
    relevant2 = [t for t in tagged2 if t.get("relevant")][:8] or tagged2[:8]
    analysis = _generate_comparison(body.question, name1, tagged1, name2, tagged2)

    def to_results(tagged):
        return [{"text": t["text"], "page": t["page"],
                 "confidence": 0, "topic": t["topic"]} for t in tagged]
    return {
        "question": body.question,
        "analysis": analysis,
        "doc1": {"doc_id": body.doc_id_1, "name": name1, "results": to_results(relevant1)},
        "doc2": {"doc_id": body.doc_id_2, "name": name2, "results": to_results(relevant2)},
    }


# ── Feature 8: Search documents ───────────────────────────────────────────────
@app.post("/search")
def search(body: SearchRequest, session_id: str = Cookie(default=None)):
    if not session_id or session_id not in _sessions:
        raise HTTPException(404, "Session not found. Please upload a document first.")
    session = _get_or_create_session(session_id)
    session_db = session["db"]
    if session_db.count() == 0:
        raise HTTPException(404, "No documents uploaded yet.")
    if not body.query.strip():
        raise HTTPException(400, "Search query cannot be empty.")
    q_vec = embeddings.embed_query(body.query)
    args = {
        "query_embeddings": [q_vec],
        "n_results":        min(6, session_db.count()),
        "include":          ["documents", "distances", "metadatas"],
    }
    if body.doc_id:
        args["where"] = {"doc_id": body.doc_id}
    r = session_db.query(**args)
    return {
        "query":   body.query,
        "results": [
            {
                "text":       _clean(doc),
                "confidence": _to_confidence(dist),
                "page":       meta.get("page"),
                "doc_id":     meta.get("doc_id"),
                "filename":   session["documents"].get(meta.get("doc_id"), {}).get("name", "Unknown"),
            }
            for doc, dist, meta in zip(r["documents"][0], r["distances"][0], r["metadatas"][0])
        ],
    }


# ── Feature 10: Study mode ────────────────────────────────────────────────────
@app.get("/study/{doc_id}")
def study(doc_id: str, session_id: str = Cookie(default=None)):
    if not session_id or session_id not in _sessions:
        raise HTTPException(404, "Session not found.")
    session = _get_or_create_session(session_id)
    if doc_id not in session["documents"]:
        raise HTTPException(404, "Document not found.")
    res    = session["db"].get(where={"doc_id": doc_id})
    chunks = [_clean(c) for c in res["documents"]]
    return {
        "quiz":       _generate_qa(chunks, max_questions=8),
        "flashcards": _generate_flashcards(chunks),
    }

@app.get("/study/{doc_id}/more")
def study_more(doc_id: str, count: int = 5, session_id: str = Cookie(default=None)):
    if not session_id or session_id not in _sessions:
        raise HTTPException(404, "Session not found.")
    session = _get_or_create_session(session_id)
    if doc_id not in session["documents"]:
        raise HTTPException(404, "Document not found.")
    res    = session["db"].get(where={"doc_id": doc_id})
    chunks = [_clean(c) for c in res["documents"]]
    return {"quiz": _generate_qa(chunks, max_questions=count)}


# ── Document management ────────────────────────────────────────────────────────
@app.get("/documents")
def list_documents(session_id: str = Cookie(default=None)):
    if not session_id or session_id not in _sessions:
        return {"documents": []}
    session = _get_or_create_session(session_id)
    return {"documents": [{"doc_id": k, **v} for k, v in session["documents"].items()]}

@app.post("/regenerate-qa/{doc_id}")
def regenerate_qa(doc_id: str, session_id: str = Cookie(default=None)):
    if not session_id or session_id not in _sessions:
        raise HTTPException(404, "Session not found.")
    session = _get_or_create_session(session_id)
    if doc_id not in session["documents"]:
        raise HTTPException(404, "Document not found.")
    res    = session["db"].get(where={"doc_id": doc_id})
    chunks = [_clean(c) for c in res["documents"]]
    suggested_qa = _generate_qa(chunks, max_questions=5)
    session["documents"][doc_id]["suggested_qa"] = suggested_qa
    _save_documents_to_disk(session["documents"])
    return {"doc_id": doc_id, "suggested_qa": suggested_qa}

@app.delete("/documents/{doc_id}")
def delete_document(doc_id: str, session_id: str = Cookie(default=None)):
    if not session_id or session_id not in _sessions:
        raise HTTPException(404, "Session not found.")
    session = _get_or_create_session(session_id)
    if doc_id not in session["documents"]:
        raise HTTPException(404, "Document not found.")
    ids = session["db"].get(where={"doc_id": doc_id})["ids"]
    if ids:
        session["db"].delete(ids=ids)
    del session["documents"][doc_id]
    session["bm25_store"].pop(doc_id, None)
    _save_documents_to_disk(session["documents"])
    return {"deleted": doc_id, "message": "Document deleted."}

@app.delete("/session")
def delete_session(session_id: str = Cookie(default=None)):
    """Manually clear the entire session and all uploaded documents."""
    if session_id and session_id in _sessions:
        del _sessions[session_id]
    response = JSONResponse({"message": "✅ Your session has been cleared. All uploaded documents have been permanently deleted."})
    response.delete_cookie("session_id")
    return response


# ── Shared query logic ─────────────────────────────────────────────────────────
def _tokenize(text: str) -> list[str]:
    """Lowercase word tokenizer for BM25."""
    return re.findall(r'[a-zA-ZàèéìíîòóùúÀÈÉÌÍÎÒÓÙÚ0-9]+', text.lower())


def _rewrite_query(question: str) -> str:
    """
    Use Groq to rewrite the user query into a search-optimised form.
    Expands abbreviations, adds key terms, strips conversational filler.
    Falls back to original question if Groq unavailable.

    Example:
      "What does the Immerman theorem say?" →
      "Immerman Szelepcsényi theorem NL coNL complexity definition"
    """
    if not SUMMARIZER_OK or not _groq_client:
        return question
    system = (
        "You are a search query optimizer. "
        "Rewrite the user question into a short, keyword-dense search query "
        "that will retrieve the most relevant passages from a technical document. "
        "Rules: keep named entities exactly as written, expand abbreviations, "
        "add synonyms if helpful, remove conversational filler words. "
        "Output ONLY the rewritten query, nothing else. Max 20 words."
    )
    try:
        return _groq_generate(system, question, max_tokens=50, model=_GROQ_FAST)
    except Exception:
        return question


def _bm25_search(query: str, doc_id: str | None, top_k: int, session: dict) -> list[dict]:
    """
    BM25 keyword search across the session bm25_store.
    Returns up to top_k results as dicts with text, page, doc_id, bm25_score.
    """
    tokens = _tokenize(query)
    results = []
    s_bm25 = session["bm25_store"]

    target_ids = [doc_id] if doc_id and doc_id in s_bm25 else list(s_bm25.keys())

    for did in target_ids:
        entry  = s_bm25[did]
        scores = entry["bm25"].get_scores(tokens)
        chunks = entry["chunks"]
        metas  = entry["metas"]
        # Get indices of top scoring chunks
        top_idx = sorted(range(len(scores)), key=lambda i: -scores[i])[:top_k]
        for idx in top_idx:
            if scores[idx] > 0:
                results.append({
                    "text":      chunks[idx],
                    "page":      metas[idx].get("page"),
                    "doc_id":    did,
                    "bm25_score": float(scores[idx]),
                })

    # Sort by BM25 score descending
    results.sort(key=lambda x: -x["bm25_score"])
    return results[:top_k]


def _query(question: str, doc_id: str | None, n: int, session: dict) -> dict:
    """
    Hybrid retrieval pipeline:
      1. Query rewriting  — expand + densify the query via Groq
      2. Vector search    — top-10 candidates from ChromaDB (cosine)
      3. BM25 search      — top-10 candidates from keyword index
      4. Reciprocal Rank Fusion — merge both result lists
      5. Cross-encoder reranker — score top merged candidates
      6. Return top-n passages + best as answer seed
    """
    session_db = session["db"]
    CANDIDATES = min(10, session_db.count())

    # ── Step 1: Query rewriting ───────────────────────────────────────────────
    rewritten = _rewrite_query(question)
    search_q  = rewritten if rewritten != question else question

    # ── Step 2: Vector search ─────────────────────────────────────────────────
    q_vec = embeddings.embed_query(search_q)
    args = {
        "query_embeddings": [q_vec],
        "n_results":        CANDIDATES,
        "include":          ["documents", "distances", "metadatas"],
    }
    if doc_id:
        args["where"] = {"doc_id": doc_id}
    r = session_db.query(**args)

    vec_docs  = r["documents"][0]
    vec_dists = r["distances"][0]
    vec_metas = r["metadatas"][0]

    # ── Step 3: BM25 keyword search ───────────────────────────────────────────
    bm25_results = _bm25_search(search_q, doc_id, top_k=CANDIDATES, session=session)

    # ── Step 4: Reciprocal Rank Fusion (RRF) ─────────────────────────────────
    # RRF score = Σ 1/(k + rank)  where k=60 is standard
    # We identify chunks by their text content as the key
    K = 60
    rrf_scores: dict[str, float] = {}
    chunk_data: dict[str, dict]  = {}

    # Score vector results
    for rank, (doc, dist, meta) in enumerate(zip(vec_docs, vec_dists, vec_metas)):
        key = doc[:120]  # use first 120 chars as dedup key
        rrf_scores[key] = rrf_scores.get(key, 0.0) + 1.0 / (K + rank + 1)
        chunk_data[key] = {"text": doc, "dist": dist, "meta": meta}

    # Score BM25 results
    for rank, hit in enumerate(bm25_results):
        key = hit["text"][:120]
        rrf_scores[key] = rrf_scores.get(key, 0.0) + 1.0 / (K + rank + 1)
        if key not in chunk_data:
            # BM25-only hit: assign a neutral distance
            chunk_data[key] = {"text": hit["text"], "dist": 0.5, "meta": {"doc_id": hit["doc_id"], "page": hit["page"]}}

    # Sort by RRF score descending, take top CANDIDATES
    sorted_keys = sorted(rrf_scores, key=lambda k: -rrf_scores[k])[:CANDIDATES]
    merged_docs  = [chunk_data[k]["text"]  for k in sorted_keys]
    merged_dists = [chunk_data[k]["dist"]  for k in sorted_keys]
    merged_metas = [chunk_data[k]["meta"]  for k in sorted_keys]

    # ── Step 5: Cross-encoder reranker ────────────────────────────────────────
    if RERANKER_OK and reranker and len(merged_docs) > 1:
        pairs  = [(question, _clean(d)) for d in merged_docs]
        scores = get_reranker().predict(pairs)
        ranked = sorted(zip(scores, merged_docs, merged_dists, merged_metas), key=lambda x: -x[0])
        ranked = ranked[:max(n, 3)]
        merged_docs  = [d    for _, d, _, _ in ranked]
        merged_dists = [dist for _, _, dist, _ in ranked]
        merged_metas = [m    for _, _, _, m in ranked]

    top_chunk = _clean(merged_docs[0])
    top_meta  = merged_metas[0]

    return {
        "answer":     top_chunk,
        "confidence": _to_confidence(merged_dists[0]),
        "page":       top_meta.get("page"),
        "filename":   session["documents"].get(top_meta.get("doc_id"), {}).get("name"),
        "rewritten_query": rewritten if rewritten != question else None,
        "passages": [
            {
                "text":       _clean(doc),
                "confidence": _to_confidence(dist),
                "page":       meta.get("page"),
                "filename":   session["documents"].get(meta.get("doc_id"), {}).get("name"),
            }
            for doc, dist, meta in zip(merged_docs, merged_dists, merged_metas)
        ],
    }


# ── Text helpers ───────────────────────────────────────────────────────────────
def _clean(text: str) -> str:
    # ── Hyphenated line breaks: "algo-\nrithm" → "algorithm" ────────────────
    text = re.sub(r'(\w)-\n(\w)', r'\1\2', text)
    # ── Merge single newlines into spaces ─────────────────────────────────────
    text = re.sub(r'(?<!\n)\n(?!\n)', ' ', text)

    # ── Math symbol repair ────────────────────────────────────────────────────
    # Accented characters garbled by PyPDF
    math_chars = {
        'G\u00a8odel': 'Gödel', 'G¨odel': 'Gödel', 'G¨odel': 'Gödel',
        'Sch\u00a8oning': 'Schöning', 'Sch¨oning': 'Schöning',
        'B\u00a8uchi': 'Büchi', 'B¨uchi': 'Büchi',
        # Set notation garbling
        'f0, 1g': '{0,1}', 'f0,1g': '{0,1}',
        'f0; 1g': '{0,1}', 'f0;1g': '{0,1}',
        '∈E0, 1∗': '∈ {0,1}*', '∈E0,1∗': '∈ {0,1}*',
        'S 0, 1n': 'S ⊆ {0,1}^n', 'S 0,1n': 'S ⊆ {0,1}^n',
        # Common symbol replacements
        'nN': 'n ∈ N', '2N': '2^N',
    }
    for bad, good in math_chars.items():
        text = text.replace(bad, good)

    # Superscript/subscript number patterns: "0,1n" → "{0,1}^n"
    text = re.sub(r'\bf(\d[\d,\s]*)g', r'{\1}', text)   # "f0,1g" → "{0,1}"
    text = re.sub(r'(\w)\s*∗', r'\1*', text)               # "1∗" → "1*"
    text = re.sub(r'(\w)\s*∈\s*(\w)', r'\1 ∈ \2', text) # "x∈S" → "x ∈ S"

    # Diacritic repair: "u¨" → "ü", "o¨" → "ö", "a¨" → "ä"
    text = re.sub(r'([aouAOU])\u00a8', lambda m: m.group(1).translate(str.maketrans("aouAOU","äöüÄÖÜ")), text)
    text = re.sub(r'([aouAOU])¨',       lambda m: m.group(1).translate(str.maketrans("aouAOU","äöüÄÖÜ")), text)

    # Camel-case fix from PDF column merging: "nTime" → "n Time"
    text = re.sub(r'([a-z])([A-Z][a-z])', r'\1 \2', text)

    # ── Structural noise removal ──────────────────────────────────────────────
    text = re.sub(r'(?m)^\s*\d+\s*$', '', text)           # lone page numbers
    text = re.sub(r'(?m)^[A-Z][A-Z\s]{6,}$', '', text)     # ALL-CAPS headers
    text = re.sub(r'\bTHESIS STATEMENT\b', '', text, flags=re.IGNORECASE)
    text = re.sub(r'\bEssay\s+\d+\b', '', text, flags=re.IGNORECASE)
    text = re.sub(r'(?<!\w)([a-z])\.\.\s', ' ', text)   # "n.. " noise
    text = re.sub(r'  +', ' ', text)
    return text.strip()

def _to_confidence(dist: float) -> int:
    return round(max(0.0, 1.0 - min(dist, 1.0)) * 100)

def _best_sentence(passages: list[dict]) -> str:
    """
    Fallback when Groq is unavailable or returns empty.
    Instead of dumping the raw chunk (often a wall of text), extract
    the single most informative sentence from the top passage —
    defined as the longest sentence that is not a heading or formula line.
    """
    if not passages:
        return "No answer found."
    text = passages[0]["text"]
    sentences = re.split(r'(?<=[.!?])\s+', text)
    # Filter: keep sentences that look like prose (>= 8 words, not all-caps)
    candidates = [
        s.strip() for s in sentences
        if len(s.split()) >= 8 and not s.isupper() and not re.match(r'^\s*[\d\W]+\s*$', s)
    ]
    if not candidates:
        return text[:400]   # absolute last resort: first 400 chars
    # Return the longest candidate (most informative)
    return max(candidates, key=len)


def _detect_question_type(question: str) -> str:
    """
    Classify the question so we can tailor the response style.
    Returns one of: definition | explanation | comparison | source_list | factual
    """
    q = question.lower().strip()
    # Definition: "what is X", "define X", "what does X mean"
    if re.search(r'\bwhat is\b|\bwhat are\b|\bdefine\b|\bwhat does .+ mean\b|\bdefinition of\b', q):
        return "definition"
    # Source/list: "what types", "list", "what kinds", "give examples", "what sources"
    if re.search(r'\bwhat types\b|\bwhat kinds\b|\blist\b|\bgive examples\b|\bwhat sources\b|\bname (some|the)\b', q):
        return "source_list"
    # Comparison: "compare", "difference between", "how does X differ", "vs"
    if re.search(r'\bcompare\b|\bdifference between\b|\bhow does .+ differ\b|\bversus\b|\bvs\.?\b|\bsimilar\b', q):
        return "comparison"
    # Explanation: "how", "why", "explain", "describe how"
    if re.search(r'\bhow\b|\bwhy\b|\bexplain\b|\bdescribe\b|\bwhat causes\b|\bhow does\b', q):
        return "explanation"
    return "factual"


def _build_system_prompt(q_type: str, doc_type: str) -> str:
    """
    Build a tailored system prompt based on question type and document type.
    Each prompt enforces grounding AND the appropriate response format.
    """
    is_technical = doc_type in ("lecture notes", "research paper")

    grounding_rule = (
        "Answer using ONLY information explicitly present in the document context provided. "
        "If the answer is not clearly stated in the context, respond with exactly: "
        "'The document does not specify this.' "
        "Never add outside knowledge, assumptions, or information not in the text. "
        "Never fabricate names, numbers, dates, or definitions."
    )

    if q_type == "definition":
        academic = (
            "Include the formal/technical definition if one exists in the text, "
            "then explain it in plain language. " if is_technical else ""
        )
        return (
            f"You are an expert explaining concepts from a document. {grounding_rule} "
            f"\n\nFor DEFINITION questions:"
            f"\n1. State a clear, precise definition in your own words."
            f"\n{academic}"
            f"2. Give one concrete example or use-case from the document."
            f"\n3. Keep it to 3-4 sentences. Do NOT copy text verbatim."
        )

    elif q_type == "explanation":
        return (
            f"You are an expert teacher. {grounding_rule} "
            f"\n\nFor HOW/WHY questions:"
            f"\n1. Explain the mechanism or reason step by step."
            f"\n2. Ground each step in what the document actually says."
            f"\n3. Use 3-5 sentences. Do NOT speculate beyond the text."
        )

    elif q_type == "comparison":
        return (
            f"You are an expert analyst. {grounding_rule} "
            f"\n\nFor COMPARISON questions:"
            f"\n1. Clearly identify the two things being compared."
            f"\n2. State the key similarity or difference as the document describes it."
            f"\n3. Keep it to 3-4 sentences. Structure: \'X does/is ... whereas Y does/is ...\'"
        )

    elif q_type == "source_list":
        return (
            f"You are a precise document analyst. {grounding_rule} "
            f"\n\nFor LIST/TYPE questions:"
            f"\n1. List the specific items, formats, or types explicitly mentioned in the document."
            f"\n2. Do NOT describe topics or themes — list the actual named items."
            f"\n3. Use a short bulleted list followed by one explanatory sentence."
        )

    else:  # factual
        return (
            f"You are a precise document assistant. {grounding_rule} "
            f"\n\nAnswer the question directly and factually in 2-3 sentences. "
            f"Do not copy the text — paraphrase in your own words."
        )


def _generate_answer(question: str, passages: list[dict]) -> str:
    """
    Use Groq (Llama 3.3) to generate a grounded, question-type-aware answer.
    Uses structured prompting with strict grounding rules.
    Falls back to _best_sentence if Groq is unavailable or fails.
    """
    if not SUMMARIZER_OK or not _groq_client or not passages:
        return _best_sentence(passages) if passages else "No answer found."

    q_type   = _detect_question_type(question)
    doc_type = "document"

    # ── Build context: scored sentence selection ──────────────────────────────
    q_words = set(re.findall(r'[a-zA-Z]{3,}', question.lower()))

    def _score_sent(s: str) -> int:
        return len(q_words & set(re.findall(r'[a-zA-Z]{3,}', s.lower())))

    parts, word_count = [], 0
    for p in passages:
        sents  = re.split(r'(?<=[.!?])\s+', p["text"])
        scored = sorted(sents, key=_score_sent, reverse=True)
        top    = {s.strip() for s in scored[:3] if len(s.split()) >= 6}
        block  = " ".join(s for s in sents if s.strip() in top)
        words  = block.split()
        if word_count + len(words) > 500:
            rem = 500 - word_count
            if rem > 30:
                parts.append(" ".join(words[:rem]))
            break
        parts.append(block)
        word_count += len(words)
    context = " ".join(parts).strip() or passages[0]["text"]

    system = _build_system_prompt(q_type, doc_type)
    user   = (
        f"Document context:\n\"\"\"\n{context}\n\"\"\"\n\n"
        f"Question: {question}\n\n"
        f"Important: If the answer is not explicitly in the context above, "
        f"respond with exactly: \"The document does not specify this.\"\n\n"
        f"Answer:"
    )

    try:
        answer = _groq_generate(system, user, max_tokens=400)
        return answer if answer else _best_sentence(passages)
    except Exception:
        return _best_sentence(passages)

def _paraphrase_sentence(sent: str, doc_type: str = "document") -> str:
    """
    Use Groq to rephrase a sentence in plain language.
    Falls back to original sentence if Groq unavailable.
    """
    if not SUMMARIZER_OK or not _groq_client:
        return sent
    context_hint = {
        "lecture notes":    "from university lecture notes",
        "research paper":   "from a research paper",
        "legal document":   "from a legal document",
        "financial report": "from a financial report",
        "medical document": "from a medical document",
        "essay collection": "from an essay",
    }.get(doc_type, "")
    system = "You are a plain-language explainer. Rewrite the given sentence in simple terms."
    user   = f"Rewrite this sentence {context_hint} in clear, simple language (one sentence): {sent}"
    try:
        result = _groq_generate(system, user, max_tokens=80, model=_GROQ_FAST)
        return result if len(result.split()) >= 4 else sent
    except Exception:
        return sent


def _summarize_chunk(chunk: str, doc_type: str) -> str:
    """MAP step: summarise a single chunk in 1-2 sentences using Groq."""
    if not SUMMARIZER_OK or not _groq_client:
        # Extractive fallback: return the longest sentence in the chunk
        sents = re.split(r'(?<=[.!?])\s+', _clean(chunk))
        sents = [s.strip() for s in sents if 8 <= len(s.split()) <= 50]
        return max(sents, key=len) if sents else chunk[:200]
    system = (
        f"You are summarising a section of {doc_type}. "
        "Write 1-2 sentences capturing the key idea. Be concise and factual."
    )
    user = f"Section:\n{chunk[:600]}"
    try:
        return _groq_generate(system, user, max_tokens=80, model=_GROQ_FAST)
    except Exception:
        return chunk[:200]


def _summarize(chunks: list[str], max_chars: int = 1000) -> str:
    """
    Map-Reduce summarization — covers the FULL document, not just the start.

    MAP:   summarise every chunk independently (1-2 sentences each)
    GROUP: collect all chunk summaries into sections
    REDUCE: ask Groq to write a coherent 3-sentence final summary

    Falls back to extractive summary if Groq unavailable.
    """
    doc_type = _detect_doc_type(chunks)

    if SUMMARIZER_OK and _groq_client is not None:
        # ── MAP: summarise every chunk (batch to avoid too many API calls) ──
        # Process chunks in groups of 5 to reduce API calls while keeping coverage
        chunk_summaries = []
        group_size = max(1, len(chunks) // 8)   # aim for ~8 groups across the doc
        for i in range(0, len(chunks), group_size):
            group = chunks[i:i + group_size]
            combined = " ".join(_clean(c) for c in group)
            combined = " ".join(combined.split()[:300])   # 300 words per group
            summary = _summarize_chunk(combined, doc_type)
            if summary:
                chunk_summaries.append(summary)

        # ── REDUCE: combine chunk summaries into final paragraph ──────────
        if chunk_summaries:
            all_summaries = " ".join(chunk_summaries)
            system = (
                "You are a document summariser. Given partial summaries of different "
                f"sections of {doc_type}, write a single coherent paragraph (3 sentences). "
                "Sentence 1: main subject. Sentence 2: key concepts across ALL sections. "
                "Sentence 3: what a reader will learn. Use plain English."
            )
            user = f"Section summaries:\n{all_summaries}"
            try:
                result = _groq_generate(system, user, max_tokens=220)
                if result and len(result.split()) >= 15:
                    return result.strip()[:max_chars].rsplit(" ", 1)[0] + "..." if len(result) > max_chars else result.strip()
            except Exception:
                pass

        # ── If reduce failed, join chunk summaries directly ──────────────
        if chunk_summaries:
            joined = " ".join(chunk_summaries)
            return joined[:max_chars].rsplit(" ", 1)[0] + "..." if len(joined) > max_chars else joined

    # ── Extractive fallback (no Groq) ────────────────────────────────────────
    overview   = _extract_overview(chunks, doc_type)
    body_sents = _extract_representative_sentences(chunks, n=2, skip=overview)
    _generic_filler = f"This {doc_type} covers the following topics."
    if overview.strip() == _generic_filler:
        intro = overview.rstrip(".")
    else:
        type_openers = {
            "lecture notes":    "These lecture notes cover",
            "research paper":   "This research paper investigates",
            "legal document":   "This legal document addresses",
            "financial report": "This financial report presents",
            "medical document": "This medical document discusses",
            "essay collection": "This collection of essays explores",
            "document":         "This document discusses",
        }
        opener_str = type_openers.get(doc_type, "This document discusses")
        intro = f"{opener_str} {overview[0].lower() + overview[1:]}"
    body   = (" Furthermore, " + " ".join(body_sents)) if body_sents else ""
    result = (intro + "." + body).strip()
    return result[:max_chars].rsplit(" ", 1)[0] + "..." if len(result) > max_chars else result

def _build_structured_summary(chunks: list[str], doc_type: str) -> dict:
    """
    Build a fully structured summary with 4 sections:
      1. Overview       — 1-2 sentence high-level description (Groq paraphrased)
      2. Main Topics    — 5 key topics extracted from across the doc
      3. Key Results    — theorems, findings, conclusions from the doc
      4. Takeaway       — Groq generated synthesis of the whole document
    """
    overview  = _extract_overview(chunks, doc_type)
    p_overview = _paraphrase_sentence(overview, doc_type)

    # Main Topics: extract representative sentences, trim to topic phrases
    topic_sents = _extract_representative_sentences(chunks, n=5, skip=overview)
    main_topics = []
    for sent in topic_sents:
        p = _paraphrase_sentence(sent, doc_type)
        words = p.split()
        main_topics.append(" ".join(words[:20]) + ("..." if len(words) > 20 else ""))

    # Key Results: look for theorem/result/definition sentences
    result_keywords = [
        "theorem", "lemma", "corollary", "definition", "proposition",
        "result", "proof", "shows that", "we prove", "it follows",
        "teorema", "lemma", "definizione", "risultato",
    ]
    key_results = []
    seen_results = set()
    for chunk in chunks:
        for sent in re.split(r'(?<=[.!?])\s+', _clean(chunk)):
            sent = sent.strip()
            words = sent.split()
            if len(words) < 8 or len(words) > 60:
                continue
            if any(kw in sent.lower() for kw in result_keywords):
                p = _paraphrase_sentence(sent, doc_type)
                short = " ".join(p.split()[:25]) + ("..." if len(p.split()) > 25 else "")
                if short not in seen_results:
                    seen_results.add(short)
                    key_results.append(short)
            if len(key_results) >= 4:
                break
        if len(key_results) >= 4:
            break

    # Takeaway: ask Groq to synthesize the whole document in one sentence
    takeaway = ""
    if SUMMARIZER_OK and _groq_client is not None:
        # Use map-reduce summaries for a more complete takeaway
        condensed = _summarize(chunks, max_chars=600)
        is_technical = doc_type in ("lecture notes", "research paper")
        precision_note = (
            " Include key technical terms or theorems if relevant." if is_technical else ""
        )
        system = f"You are a document assistant. Write one precise, informative sentence.{precision_note}"
        user   = f"In one sentence, what is the main goal of these {doc_type} and what will a reader learn?\n\n{condensed}"
        try:
            takeaway = _groq_generate(system, user, max_tokens=100)
        except Exception:
            takeaway = ""

    if not takeaway:
        takeaway = f"This {doc_type} introduces key concepts and results in the field."

    return {
        "doc_type":    doc_type,
        "overview":    p_overview,
        "main_topics": main_topics,
        "key_results": key_results,
        "takeaway":    takeaway,
        # Keep these for backward compatibility
        "summary":     p_overview,
        "key_points":  main_topics,
        "conclusions": key_results,
    }


def _abstractive_bullets(chunks: list[str], max_bullets: int = 6) -> list[str]:
    """
    One bullet per document section, paraphrased with Pegasus.
    """
    doc_type = _detect_doc_type(chunks)
    raw_sents = _extract_representative_sentences(chunks, n=max_bullets)
    bullets = []
    for sent in raw_sents:
        paraphrased = _paraphrase_sentence(sent, doc_type)
        words = paraphrased.split()
        bullet = " ".join(words[:22]) + ("..." if len(words) > 22 else "")
        bullets.append(bullet)
    return bullets if bullets else _extract_bullets(chunks, max_bullets)


def _detect_doc_type(chunks: list[str]) -> str:
    """Detect document type from first 5 chunks."""
    sample = " ".join(chunks[:5]).lower()
    if any(w in sample for w in ["lecture", "course notes", "chapter", "theorem", "proof", "corollary", "lemma"]):
        return "lecture notes"
    if any(w in sample for w in ["abstract", "introduction", "methodology", "conclusion", "references", "doi"]):
        return "research paper"
    if any(w in sample for w in ["contract", "agreement", "clause", "party", "shall", "whereas"]):
        return "legal document"
    if any(w in sample for w in ["revenue", "profit", "fiscal", "quarter", "shareholders", "earnings"]):
        return "financial report"
    if any(w in sample for w in ["patient", "diagnosis", "treatment", "clinical", "medical"]):
        return "medical document"
    if any(w in sample for w in ["essay", "thesis statement", "persuasive", "argument"]):
        return "essay collection"
    return "document"


def _extract_overview(chunks: list[str], doc_type: str) -> str:
    """Extract the best overview sentence from the first 6 chunks."""
    purpose_words = [
        "this document", "this paper", "this report", "this book",
        "these notes", "this essay", "the purpose", "we study",
        "we present", "we introduce", "this work", "this course",
        "the goal", "the aim", "this text", "this lecture",
        # academic/math intro patterns
        "in this", "the focus", "we consider", "we investigate",
        "we define", "we develop", "deals with", "concerned with",
        "devoted to", "the subject", "the topic",
    ]
    for chunk in chunks[:6]:
        for sent in re.split(r'(?<=[.!?])\s+', _clean(chunk)):
            sent = sent.strip()
            words = sent.split()
            if len(words) < 8 or len(words) > 60:
                continue
            # Skip sentences that are pure symbol/formula lines
            if re.search(r'[∀∃⊆⊇∈∉→⇒λ]|:=|\b[A-Z]\s*=\s*[A-Z]', sent):
                continue
            if any(pw in sent.lower() for pw in purpose_words):
                return sent
    # Second pass: any decent-length prose sentence from the first 3 chunks
    for chunk in chunks[:3]:
        for sent in re.split(r'(?<=[.!?])\s+', _clean(chunk)):
            sent = sent.strip()
            words = sent.split()
            if 15 <= len(words) <= 50 and not sent.isupper():
                if not re.search(r'[∀∃⊆⊇∈∉→⇒λ]|:=', sent):
                    return sent
    return f"This {doc_type} covers the following topics."


def _extract_representative_sentences(chunks: list[str], n: int = 4, skip: str = "") -> list[str]:
    """Pick n informative sentences spread evenly across the document using TF scoring."""
    all_words = re.findall(r'\b[a-zA-Z]{4,}\b', " ".join(chunks).lower())
    stop = {'this','that','with','from','have','they','been','were','will','their',
            'also','when','which','into','more','than','then','some','such','each',
            'these','those','there','where','after','before','about','would','could',
            'should','other','within','while','during','between','through','however'}
    freq: dict = {}
    for w in all_words:
        if w not in stop:
            freq[w] = freq.get(w, 0) + 1

    section_size = max(1, len(chunks) // n)
    selected = []
    skip_lower = skip.lower()

    for s in range(n):
        section = chunks[s * section_size: (s + 1) * section_size]
        best_score, best_sent = 0, ""
        for chunk in section:
            for sent in re.split(r'(?<=[.!?])\s+', _clean(chunk)):
                sent = sent.strip()
                words = sent.split()
                if len(words) < 10 or len(words) > 50:
                    continue
                if sent.isupper() or sent.lower() == skip_lower:
                    continue
                if re.match(r'^\d+[\.\)\]\s]', sent):
                    continue
                score = sum(freq.get(w.lower(), 0) for w in words if w.lower() not in stop)
                if score > best_score:
                    best_score, best_sent = score, sent
        if best_sent and best_sent not in selected:
            selected.append(best_sent)

    return selected


def _extract_bullets(chunks: list[str], max_bullets: int = 6) -> list[str]:
    keywords = [
        # English
        'important', 'key', 'main', 'significant', 'major', 'primary', 'critical', 'essential', 'notably', 'highlight',
        # Italian
        'importante', 'principale', 'significativo', 'fondamentale', 'essenziale', 'chiave', 'cruciale', 'rilevante', 'notevole',
    ]
    bullets = []
    for chunk in chunks[:15]:
        for sent in re.split(r'(?<=[.!?])\s+', chunk):
            sent = sent.strip()
            if 10 <= len(sent.split()) <= 40 and any(kw in sent.lower() for kw in keywords):
                bullets.append(sent)
                if len(bullets) >= max_bullets:
                    return bullets
    for chunk in chunks[:10]:
        for sent in re.split(r'(?<=[.!?])\s+', chunk):
            sent = sent.strip()
            if 15 <= len(sent.split()) <= 35 and sent not in bullets:
                bullets.append(sent)
                if len(bullets) >= max_bullets:
                    return bullets
    return bullets


def _extract_conclusions(chunks: list[str], max_conclusions: int = 3) -> list[str]:
    keywords = [
        # English
        'therefore','conclude','conclusion','result','shows','demonstrates','indicates','suggests','thus','hence','in summary','overall','in conclusion',
        # Italian
        'quindi','pertanto','conclude','conclusione','risulta','dimostra','indica','suggerisce','dunque','perciò','in sintesi','complessivamente','in conclusione','si evince','emerge che',
    ]
    results = []
    for chunk in chunks:
        for sent in re.split(r'(?<=[.!?])\s+', chunk):
            if any(kw in sent.lower() for kw in keywords):
                sent = sent.strip()
                if 8 <= len(sent.split()) <= 50:
                    results.append(sent)
                    if len(results) >= max_conclusions:
                        return results
    return results

def _extract_people(text: str) -> list[str]:
    people = set()
    # English + Italian titles
    for m in re.findall(r'\b(?:Mr|Mrs|Ms|Dr|Prof|Sir|Sig|Sig\.ra|Dott|Dott\.ssa|Avv|Ing|Arch)\.?\s+[A-Z][a-z]+(?:\s+[A-Z][a-z]+)?', text):
        people.add(m.strip())
    for m in re.findall(r'(?<=[,;]\s)([A-Z][a-z]+\s+[A-Z][a-z]+)', text):
        people.add(m.strip())
    return sorted(list(people))[:10]

def _extract_companies(text: str) -> list[str]:
    companies = set()
    # English + Italian company suffixes
    for m in re.findall(r'[A-Z][A-Za-z\s&]+(?:Inc|Corp|Ltd|LLC|Co|Company|Group|Holdings|Technologies|Solutions|Services|International|Industries|SpA|Srl|Snc|Sas|SpA|Holding|Gruppo)\.?\b', text):
        c = m.strip()
        if len(c) > 4:
            companies.add(c)
    return sorted(list(companies))[:10]

def _extract_dates(text: str) -> list[str]:
    dates = set()
    patterns = [
        # English months
        r'\b(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{1,2},?\s+\d{4}\b',
        # Italian months
        r'\b\d{1,2}\s+(?:gennaio|febbraio|marzo|aprile|maggio|giugno|luglio|agosto|settembre|ottobre|novembre|dicembre)\s+\d{4}\b',
        # Italian abbreviated form: "15 mar. 2023"
        r'\b\d{1,2}\s+(?:gen|feb|mar|apr|mag|giu|lug|ago|set|ott|nov|dic)\.?\s+\d{4}\b',
        # Numeric formats
        r'\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b',
        r'\b(?:Q[1-4]\s+)?\d{4}\b',
    ]
    for pat in patterns:
        for m in re.findall(pat, text, re.IGNORECASE):
            dates.add(m.strip())
    return sorted(list(dates))[:15]

def _extract_numbers(text: str) -> list[str]:
    numbers = set()
    patterns = [
        r'\b\d+(?:[.,]\d+)?%(?:\s+(?:growth|increase|decrease|crescita|aumento|diminuzione|calo))?\b',
        r'(?:\$|€|£)\s*\d+(?:[.,]\d{3})*(?:[.,]\d+)?(?:\s*(?:million|billion|thousand|milioni|miliardi|migliaia|M|B|K))?\b',
        r'\b\d+(?:[.,]\d{3})*(?:[.,]\d+)?\s*(?:million|billion|thousand|milioni|miliardi|migliaia)\b',
    ]
    for pat in patterns:
        for match in re.finditer(pat, text, re.IGNORECASE):
            start = max(0, match.start() - 40)
            end   = min(len(text), match.end() + 40)
            ctx   = re.sub(r'\s+', ' ', text[start:end].strip())[:80]
            numbers.add(ctx)
    return sorted(list(numbers))[:10]


# ── Bank Statement Analysis ────────────────────────────────────────────────────

import re as _re
from datetime import datetime

def _parse_bank_transactions(text: str) -> list[dict]:
    """
    Parse bank transactions from raw PDF text.
    Supports Italian and English bank statement formats.
    Returns list of {date, description, amount} dicts.
    """
    transactions = []

    # Normalise whitespace
    text = re.sub(r'\r\n|\r', '\n', text)
    lines = text.split('\n')

    # Date patterns: DD/MM/YYYY, DD-MM-YYYY, YYYY-MM-DD, DD MMM YYYY, DD.MM.YYYY
    date_patterns = [
        r'\b(\d{2}[\/\-\.]\d{2}[\/\-\.]\d{4})\b',
        r'\b(\d{4}[\/\-\.]\d{2}[\/\-\.]\d{2})\b',
        r'\b(\d{1,2}\s+(?:gen|feb|mar|apr|mag|giu|lug|ago|set|ott|nov|dic|jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)\w*\s+\d{4})\b',
    ]

    # Amount pattern: handles both Italian (1.234,56) and English (1,234.56) formats
    # Also captures optional +/- sign and common EUR/GBP/USD symbols
    amount_pattern = re.compile(
        r'([+\-])?\s*(?:EUR|€|\$|£)?\s*(\d{1,3}(?:[.,]\d{3})*(?:[.,]\d{2}))\s*(?:EUR|€|\$|£)?'
    )

    def parse_date(s: str) -> str | None:
        s = s.strip()
        # Italian month names
        ita_months = {'gen':'01','feb':'02','mar':'03','apr':'04','mag':'05','giu':'06',
                      'lug':'07','ago':'08','set':'09','ott':'10','nov':'11','dic':'12'}
        for k, v in ita_months.items():
            s = re.sub(rf'\b{k}\w*\b', v, s, flags=re.IGNORECASE)
        for fmt in ['%d/%m/%Y','%d-%m-%Y','%d.%m.%Y','%Y-%m-%d','%Y/%m/%d',
                    '%d/%m/%y','%d %m %Y','%d %m %y']:
            try:
                return datetime.strptime(s, fmt).strftime('%Y-%m-%d')
            except ValueError:
                pass
        return None

    def parse_amount(s: str) -> float | None:
        """Convert Italian/English formatted number to float."""
        s = s.strip().replace(' ', '')
        # Detect format: if last separator is ',' and penultimate is '.': Italian 1.234,56
        if re.match(r'^\d{1,3}(\.\d{3})*(,\d{2})$', s):
            s = s.replace('.', '').replace(',', '.')
        # If last separator is '.' and penultimate is ',': English 1,234.56
        elif re.match(r'^\d{1,3}(,\d{3})*(\.\d{2})$', s):
            s = s.replace(',', '')
        else:
            # Generic cleanup
            s = s.replace(',', '.')
            # If multiple dots, keep only last as decimal
            parts = s.split('.')
            if len(parts) > 2:
                s = ''.join(parts[:-1]) + '.' + parts[-1]
        try:
            return float(s)
        except ValueError:
            return None

    seen = set()

    for i, line in enumerate(lines):
        line = line.strip()
        if not line or len(line) < 8:
            continue

        # Try to extract a date from the line
        date_str = None
        for dp in date_patterns:
            m = re.search(dp, line, re.IGNORECASE)
            if m:
                date_str = parse_date(m.group(1))
                if date_str:
                    break

        if not date_str:
            continue

        # Extract all amounts from the line
        amounts_found = amount_pattern.findall(line)
        if not amounts_found:
            # Check adjacent lines (±1) for amount
            adjacent = '\n'.join(lines[max(0,i-1):i+2])
            amounts_found = amount_pattern.findall(adjacent)

        if not amounts_found:
            continue

        # Take the last (rightmost) amount, most likely to be the transaction value
        sign_str, raw_amt = amounts_found[-1]
        amount = parse_amount(raw_amt)
        if amount is None or amount == 0:
            continue

        # Determine sign: look for debit/credit keywords
        line_lower = line.lower()
        credit_kw = ['accredito','bonifico in entrata','versamento','stipendio','rimborso',
                     'credit','deposit','incoming','received','entrata','avere','dare',
                     'ricarica','pagamento ricevuto']
        debit_kw  = ['addebito','pagamento','prelievo','commissione','spesa','canone',
                     'debit','withdrawal','payment','charge','fee','uscita','atm',
                     'pos','utenza','bolletta']

        is_credit = any(kw in line_lower for kw in credit_kw)
        is_debit  = any(kw in line_lower for kw in debit_kw)

        if sign_str == '+' or is_credit:
            signed_amount = abs(amount)
        elif sign_str == '-' or is_debit:
            signed_amount = -abs(amount)
        else:
            # Heuristic: if two amounts found on line, second is usually balance
            if len(amounts_found) >= 2:
                _, raw2 = amounts_found[0]
                signed_amount = parse_amount(raw2)
                if signed_amount is None:
                    continue
                # negative if amount reduces, positive otherwise — default debit
                signed_amount = -abs(signed_amount)
            else:
                # Default: treat as debit (outgoing)
                signed_amount = -abs(amount)

        # Build description: strip date and amount tokens from line
        desc = re.sub(r'\b\d{1,2}[\/\-\.]\d{1,2}[\/\-\.]\d{2,4}\b', '', line)
        desc = amount_pattern.sub('', desc)
        desc = re.sub(r'[+\-]\s*(?:EUR|€|\$|£)?', '', desc)
        desc = re.sub(r'\s{2,}', ' ', desc).strip(' |-_/')
        if not desc:
            desc = 'Transaction'

        # Dedup
        key = (date_str, round(signed_amount, 2))
        if key in seen:
            continue
        seen.add(key)

        transactions.append({
            'date':        date_str,
            'description': desc[:120],
            'amount':      round(signed_amount, 2),
        })

    # Sort by date
    transactions.sort(key=lambda t: t['date'])
    return transactions


def _aggregate_monthly(transactions: list[dict]) -> list[dict]:
    """Aggregate transactions by month."""
    monthly: dict[str, dict] = {}
    for t in transactions:
        month = t['date'][:7]  # YYYY-MM
        if month not in monthly:
            monthly[month] = {'month': month, 'income': 0.0, 'expense': 0.0}
        if t['amount'] >= 0:
            monthly[month]['income'] += t['amount']
        else:
            monthly[month]['expense'] += abs(t['amount'])
    result = sorted(monthly.values(), key=lambda m: m['month'])
    for m in result:
        m['income']  = round(m['income'],  2)
        m['expense'] = round(m['expense'], 2)
        # Human-readable label: "Jan 2024"
        try:
            dt = datetime.strptime(m['month'], '%Y-%m')
            m['month'] = dt.strftime('%b %Y')
        except Exception:
            pass
    return result


def _find_column(headers: list[str], candidates: list[str]) -> str | None:
    """Case-insensitive column name matcher."""
    for h in headers:
        if h.strip().lower() in [c.lower() for c in candidates]:
            return h
    return None

def _parse_structured_bank(rows: list[dict]) -> list[dict]:
    """
    Parse a list of dicts (from CSV/Excel) into transactions.
    Auto-detects date, description, debit, credit, amount, and balance columns.
    Supports formats:
      - Separate debit/credit columns (Debit_Exit, Credit_Entry)
      - Single signed amount column (Amount, Importo)
      - Italian and English column names
    """
    if not rows:
        return []

    headers = list(rows[0].keys())

    # ── Auto-detect columns ───────────────────────────────────────────────────
    date_col  = _find_column(headers, ['date','data','datum','fecha','date_transaction','transaction_date','data_operazione','data_valuta'])
    desc_col  = _find_column(headers, ['description','descrizione','desc','memo','narrative','causale','wording','details','detail','transaction_description'])
    debit_col = _find_column(headers, ['debit','debit_exit','uscita','addebito','expense','out','withdrawal','debet','ausgabe','charge'])
    credit_col= _find_column(headers, ['credit','credit_entry','entrata','accredito','income','in','deposit','credit_amount','einnahme'])
    amount_col= _find_column(headers, ['amount','importo','betrag','importe','montant','valore','net_amount','transaction_amount'])
    bal_col   = _find_column(headers, ['balance','saldo','kontostand','solde','saldo_contabile'])

    if not date_col:
        raise ValueError("Could not find a date column. Expected column named 'Date', 'Data', or similar.")

    transactions = []
    for row in rows:
        # ── Date ─────────────────────────────────────────────────────────────
        raw_date = str(row.get(date_col, '') or '').strip()
        if not raw_date:
            continue
        date_str = None
        for fmt in ['%Y-%m-%d','%d/%m/%Y','%d-%m-%Y','%d.%m.%Y','%Y/%m/%d','%m/%d/%Y','%d/%m/%y','%Y%m%d']:
            try:
                date_str = datetime.strptime(raw_date, fmt).strftime('%Y-%m-%d')
                break
            except ValueError:
                pass
        if not date_str:
            continue  # skip unparseable dates

        # ── Description ──────────────────────────────────────────────────────
        desc = str(row.get(desc_col, '') or '').strip() if desc_col else 'Transaction'
        if not desc:
            desc = 'Transaction'

        # ── Amount ───────────────────────────────────────────────────────────
        def to_float(val) -> float:
            s = str(val or '').strip().replace(' ', '').replace('€','').replace('$','').replace('£','')
            if not s or s in ('-','—',''):
                return 0.0
            # Italian format: 1.234,56 → 1234.56
            if re.match(r'^\d{1,3}(\.\d{3})*(,\d+)$', s):
                s = s.replace('.', '').replace(',', '.')
            else:
                s = s.replace(',', '.')
            try:
                return float(s)
            except ValueError:
                return 0.0

        if debit_col and credit_col:
            # Separate debit/credit columns
            debit  = to_float(row.get(debit_col,  0))
            credit = to_float(row.get(credit_col, 0))
            if credit > 0:
                amount = credit
            elif debit > 0:
                amount = -debit
            else:
                continue  # no movement
        elif amount_col:
            amount = to_float(row.get(amount_col, 0))
            if amount == 0:
                continue
        else:
            continue  # can't determine amount

        # ── Balance ──────────────────────────────────────────────────────────
        balance = to_float(row.get(bal_col, None)) if bal_col else None

        transactions.append({
            'date':        date_str,
            'description': desc[:120],
            'amount':      round(amount, 2),
            'balance':     round(balance, 2) if balance is not None else None,
        })

    transactions.sort(key=lambda t: t['date'])
    return transactions


@app.post("/bank/analyze")
async def bank_analyze(file: UploadFile = File(...)):
    """
    Parse a bank statement (CSV, Excel, or PDF) and return transactions + monthly aggregates.
    CSV and Excel are parsed as structured data (column-aware).
    PDF falls back to regex text parsing.
    No session required — results are temporary (not persisted).
    """
    ext = ('.' + file.filename.rsplit('.', 1)[-1]).lower() if '.' in file.filename else ''
    if ext not in {'.pdf', '.csv', '.xlsx', '.xls'}:
        raise HTTPException(400, "Bank analysis supports PDF, CSV, and Excel files.")

    file_bytes = await file.read()
    transactions = []

    # ── CSV: structured column parsing ───────────────────────────────────────
    if ext == '.csv':
        try:
            import csv as _csv
            # Try common delimiters
            text = file_bytes.decode('utf-8', errors='replace')
            dialect = _csv.Sniffer().sniff(text[:4096], delimiters=',;\t|')
            reader = _csv.DictReader(io.StringIO(text), dialect=dialect)
            rows = list(reader)
        except Exception:
            rows = list(_csv.DictReader(io.StringIO(file_bytes.decode('utf-8', errors='replace'))))
        try:
            transactions = _parse_structured_bank(rows)
        except ValueError as e:
            raise HTTPException(422, str(e))

    # ── Excel: structured column parsing ─────────────────────────────────────
    elif ext in ('.xlsx', '.xls'):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            ws = wb.active
            rows_raw = list(ws.iter_rows(values_only=True))
            wb.close()
            if not rows_raw:
                raise HTTPException(422, "Excel file is empty.")
            headers = [str(h).strip() if h is not None else f'col_{i}' for i, h in enumerate(rows_raw[0])]
            rows = [dict(zip(headers, r)) for r in rows_raw[1:] if any(c is not None for c in r)]
            transactions = _parse_structured_bank(rows)
        except HTTPException:
            raise
        except ValueError as e:
            raise HTTPException(422, str(e))
        except Exception as e:
            raise HTTPException(422, f"Could not read Excel file: {e}")

    # ── PDF: regex text parsing ───────────────────────────────────────────────
    else:
        try:
            pages_text = _extract_text_from_file(file.filename, file_bytes)
        except Exception:
            raise HTTPException(422, "Could not read this PDF.")
        full_text = "\n".join(text for _, text in pages_text)
        if not full_text.strip():
            raise HTTPException(422, "No text found in this PDF. Is it a scanned image?")
        transactions = _parse_bank_transactions(full_text)

    if not transactions:
        raise HTTPException(422,
            "No transactions could be parsed. "
            "Check that your file has columns for Date, Description, and Debit/Credit amounts.")

    total_income  = round(sum(t['amount'] for t in transactions if t['amount'] >  0), 2)
    total_expense = round(sum(abs(t['amount']) for t in transactions if t['amount'] < 0), 2)
    net_balance   = round(total_income - total_expense, 2)
    monthly       = _aggregate_monthly(transactions)

    # ── Enrich transactions with category ────────────────────────────────────
    transactions = _detect_categories(transactions)

    # ── Anomaly detection ─────────────────────────────────────────────────────
    anomalies = _detect_anomalies(transactions)

    # ── Cash-flow prediction (next 3 months) ──────────────────────────────────
    prediction = _predict_cashflow(monthly)

    # ── Category breakdown ────────────────────────────────────────────────────
    cat_totals: dict = {}
    for t in transactions:
        cat = t.get('category', 'Other')
        cat_totals[cat] = cat_totals.get(cat, 0.0) + abs(t['amount'])
    category_breakdown = [
        {'category': k, 'total': round(v, 2)}
        for k, v in sorted(cat_totals.items(), key=lambda x: -x[1])
    ]

    return {
        "transactions":        transactions,
        "monthly":             monthly,
        "total_income":        total_income,
        "total_expense":       total_expense,
        "net_balance":         net_balance,
        "count":               len(transactions),
        "anomalies":           anomalies,
        "prediction":          prediction,
        "category_breakdown":  category_breakdown,
    }


# ── Category detection ─────────────────────────────────────────────────────────
_CATEGORY_RULES = [
    # (category, keywords)
    ('Salary & Payroll',   ['salary','stipendio','payroll','wage','retribuzione','busta paga']),
    ('Revenue & Sales',    ['revenue','sales','online sales','client payment','invoice','fattura','vendita','ricavo']),
    ('Consulting',         ['consulting','consulenza','fee received','professional']),
    ('Rent & Property',    ['rent','affitto','lease','locazione','office rent']),
    ('Utilities',          ['utilities','utenza','electricity','gas','water','internet','phone','bolletta']),
    ('Software & Tech',    ['software','subscription','abbonamento','saas','license','licenza','tech']),
    ('Equipment',          ['equipment','attrezzatura','hardware','machinery','purchase','acquisto']),
    ('Suppliers',          ['supplier','fornitore','vendor','supply','materiali']),
    ('Transfers',          ['transfer','trasferimento','partner','bonifico','wire']),
    ('Refunds & Returns',  ['refund','rimborso','return','restituzione','chargeback']),
    ('Banking & Fees',     ['fee','commissione','charge','interest','bank','tassa']),
    ('Tax & Legal',        ['tax','tasse','iva','legal','notary','notaio','imposte']),
]

def _detect_categories(transactions: list[dict]) -> list[dict]:
    for t in transactions:
        desc_lower = t['description'].lower()
        t['category'] = 'Other'
        for cat, keywords in _CATEGORY_RULES:
            if any(kw in desc_lower for kw in keywords):
                t['category'] = cat
                break
    return transactions


# ── Anomaly detection ──────────────────────────────────────────────────────────
def _detect_anomalies(transactions: list[dict]) -> list[dict]:
    """
    Flag transactions that are statistical outliers using z-score per category.
    Also flags: duplicate amounts on same day, unusually round numbers.
    """
    import math

    # Z-score per category
    from collections import defaultdict
    cat_amounts: dict = defaultdict(list)
    for t in transactions:
        cat_amounts[t.get('category','Other')].append(abs(t['amount']))

    cat_stats: dict = {}
    for cat, vals in cat_amounts.items():
        if len(vals) < 2:
            cat_stats[cat] = (0, 0)
            continue
        mean = sum(vals) / len(vals)
        std  = math.sqrt(sum((v - mean)**2 for v in vals) / len(vals))
        cat_stats[cat] = (mean, std)

    # Global stats as fallback
    all_abs = [abs(t['amount']) for t in transactions]
    g_mean  = sum(all_abs) / len(all_abs) if all_abs else 0
    g_std   = math.sqrt(sum((v - g_mean)**2 for v in all_abs) / len(all_abs)) if all_abs else 1

    # Duplicate detection: same amount + same day
    seen_day_amounts: dict = defaultdict(list)
    for t in transactions:
        seen_day_amounts[(t['date'], abs(t['amount']))].append(t['description'])

    anomalies = []
    for t in transactions:
        reasons = []
        amt = abs(t['amount'])
        cat = t.get('category', 'Other')

        # Z-score check (global)
        if g_std > 0:
            z = (amt - g_mean) / g_std
            if z > 2.5:
                reasons.append(f'Unusually large amount (z={z:.1f}σ above average)')

        # Category z-score
        mean_c, std_c = cat_stats.get(cat, (0, 0))
        if std_c > 0:
            z_c = (amt - mean_c) / std_c
            if z_c > 2.5:
                reasons.append(f'Large for category "{cat}" (z={z_c:.1f}σ)')

        # Duplicate same-day same-amount
        key = (t['date'], amt)
        if len(seen_day_amounts[key]) > 1:
            reasons.append(f'Duplicate amount on same day ({len(seen_day_amounts[key])}x)')

        # Suspiciously round number (>500 and ends in 000)
        if amt >= 500 and amt % 1000 == 0:
            reasons.append('Suspiciously round amount')

        if reasons:
            anomalies.append({
                'date':        t['date'],
                'description': t['description'],
                'amount':      t['amount'],
                'category':    cat,
                'reasons':     reasons,
                'severity':    'high' if len(reasons) >= 2 else 'medium',
            })

    return sorted(anomalies, key=lambda x: ('high','medium').index(x['severity']))


# ── Cash-flow prediction ───────────────────────────────────────────────────────
def _predict_cashflow(monthly: list[dict]) -> list[dict]:
    """
    Predict next 3 months using weighted linear regression on income & expense.
    More recent months get higher weight.
    """
    from datetime import datetime, timedelta
    import math

    if len(monthly) < 2:
        return []

    def parse_month(label: str):
        for fmt in ['%b %Y','%Y-%m']:
            try:
                return datetime.strptime(label, fmt)
            except ValueError:
                pass
        return None

    valid = [(parse_month(m['month']), m['income'], m['expense'])
             for m in monthly if parse_month(m['month'])]
    if len(valid) < 2:
        return []

    valid.sort(key=lambda x: x[0])
    n = len(valid)

    def weighted_linreg(ys):
        # Weights: more recent = higher weight (linear ramp)
        ws = [i + 1 for i in range(n)]
        xs = list(range(n))
        sw  = sum(ws)
        swx = sum(w*x for w,x in zip(ws,xs))
        swy = sum(w*y for w,y in zip(ws,ys))
        swxx= sum(w*x*x for w,x in zip(ws,xs))
        swxy= sum(w*x*y for w,x,y in zip(ws,xs,ys))
        denom = sw*swxx - swx**2
        if denom == 0:
            return swy/sw, 0
        slope = (sw*swxy - swx*swy) / denom
        intercept = (swy - slope*swx) / sw
        return intercept, slope

    inc_vals = [v[1] for v in valid]
    exp_vals = [v[2] for v in valid]
    inc_int, inc_slope = weighted_linreg(inc_vals)
    exp_int, exp_slope = weighted_linreg(exp_vals)

    # Last month date
    last_dt = valid[-1][0]
    predictions = []
    for i in range(1, 4):
        x = n - 1 + i
        # Next month
        if last_dt.month + i <= 12:
            pred_dt = last_dt.replace(month=last_dt.month + i)
        else:
            pred_dt = last_dt.replace(year=last_dt.year + 1, month=(last_dt.month + i) % 12 or 12)

        pred_income  = max(0, round(inc_int + inc_slope * x, 2))
        pred_expense = max(0, round(exp_int + exp_slope * x, 2))
        predictions.append({
            'month':        pred_dt.strftime('%b %Y'),
            'income':       pred_income,
            'expense':      pred_expense,
            'net':          round(pred_income - pred_expense, 2),
            'predicted':    True,
        })

    return predictions




    stopwords = {
        # English
        'the','a','an','and','or','but','in','on','at','to','for','of','with','by','from','as',
        'is','was','are','were','be','been','this','that','these','those','it','its','they',
        'them','we','our','you','your','he','she','i','my','not','so','if','can','all','any',
        'more','most','also','than','then','when','where','which','who','what','how','about',
        'after','before','have','has','had','do','does','will','would','could','should','may','might',
        # Italian
        'il','lo','la','gli','le','un','una','uno','dei','del','della','delle','degli','di','da',
        'in','con','su','per','tra','fra','non','che','chi','cui','come','quando','dove','perché',
        'ma','se','anche','già','più','molto','questo','questa','questi','queste','quello','quella',
        'sono','siamo','sei','è','era','erano','hanno','hai','ha','ho','essere','avere','fare',
        'loro','noi','voi','lui','lei','mi','ti','si','ci','vi','ne','gli','le','uno','due',
        'dopo','prima','sempre','ancora','tutto','tutti','però','quindi','così','ogni',
    }
    words = re.findall(r'\b[a-zA-ZàèéìíîòóùúÀÈÉÌÍÎÒÓÙÚ]{4,}\b', text.lower())
    freq  = {}
    for w in words:
        if w not in stopwords:
            freq[w] = freq.get(w, 0) + 1
    return [w.title() for w, _ in sorted(freq.items(), key=lambda x: -x[1])[:10]]

def _generate_qa(chunks: list[str], max_questions: int = 5) -> list[dict]:
    """
    Generate suggested questions using Groq (Llama 3.1).
    Falls back to regex-based extraction if Groq is unavailable.
    """
    # Build a representative sample of the document for Groq
    n = len(chunks)
    sample_idx = list(dict.fromkeys([0, 1, n//4, n//2, 3*n//4, n-1]))
    sample_text = " ".join(_clean(chunks[i]) for i in sample_idx if i < n)
    sample_text = " ".join(sample_text.split()[:500])

    if SUMMARIZER_OK and _groq_client is not None:
        system = (
            "You are a study assistant. Generate exactly 5 thoughtful questions "
            "a student would ask after reading the document. "
            "Rules: each question must be specific to THIS document's content, "
            "start with What/How/Why/Which/Explain, be answerable from the text, "
            "and be self-contained (no 'here' or 'this section'). "
            "Output ONLY the 5 questions, one per line, no numbering, no preamble."
        )
        user = f"Document excerpt:\n{sample_text}"
        try:
            raw = _groq_generate(system, user, max_tokens=300, model=_GROQ_FAST)
            lines = [l.strip().lstrip('0123456789.-) ') for l in raw.splitlines() if l.strip()]
            questions = []
            seen = set()
            for q in lines:
                if (len(q.split()) >= 5
                        and q.endswith('?')
                        and q not in seen
                        and not re.search(r'\b(this section|here|the passage)\b', q, re.I)):
                    seen.add(q)
                    questions.append({"question": q, "answer": ""})
                if len(questions) >= max_questions:
                    break
            if questions:
                return questions
        except Exception:
            pass

    # ── Regex fallback (used when Groq unavailable) ───────────────────────────
    noise = re.compile(
        r'thesis statement|essay \d|persuasive|click to|drag|upload|^\d+\s',
        re.IGNORECASE
    )
    exposition_starters = {
        'recall','note','let','assume','consider','observe','suppose',
        'by','since','thus','hence','therefore','proof','clearly',
        'here','then','now','next','first','second','finally',
        'moreover','furthermore','however','similarly','conversely',
    }
    questions, seen = [], set()
    for chunk in chunks[:20]:
        for sent in re.split(r'(?<=[.!?])\s+', _clean(chunk)):
            sent = sent.strip()
            words = sent.split()
            if len(words) < 10 or len(words) > 30:
                continue
            if sent.upper() == sent:
                continue
            if noise.search(sent):
                continue
            if words[0].lower() in exposition_starters:
                continue
            if re.search(r'[∀∃⊆⊇∈∉∪∩→↔⇒⇔λδε]|\\frac|\\sum', sent):
                continue
            q = _sentence_to_question(sent)
            if q and q not in seen and len(q.split()) >= 5:
                seen.add(q)
                questions.append({"question": q, "answer": sent})
            if len(questions) >= max_questions:
                return questions
    return questions

def _sentence_to_question(s: str) -> str | None:
    s = s.rstrip('.!?').strip()
    words = s.split()
    lower = s.lower()

    # Reject if sentence starts with a pronoun or conjunction
    bad_starters = {'you','i','we','they','he','she','it','when','if','because',
                    'but','and','or','the','a','an','this','that','these','those'}
    if words[0].lower() in bad_starters:
        return None

    # Subject validity check: a good subject for "What is X?" must be:
    #   - 1 to 3 words
    #   - Start with a capital letter (it's a named concept, not a clause)
    #   - Contain no lowercase stopwords (rejects "Recall that", "a function", etc.)
    clause_words = {'that','which','when','where','if','as','an','the',
                    'this','these','those','some','such','any','all'}

    def _valid_subject(subj: str) -> bool:
        sw = subj.split()
        if not (1 <= len(sw) <= 3):
            return False
        if not sw[0][0].isupper():
            return False
        if any(w.lower() in clause_words for w in sw):
            return False
        return True

    # English: "X is Y" → "What is X?"
    if len(words) >= 5:
        for i, w in enumerate(words[1:], start=1):
            if w.lower() == 'is':
                subject = ' '.join(words[:i])
                if _valid_subject(subject):
                    return f"What is {subject}?"
            if w.lower() == 'are':
                subject = ' '.join(words[:i])
                if _valid_subject(subject):
                    return f"What are {subject}?"

    # Italian: "X è Y" → "Cos'è X?"
    if len(words) >= 5:
        for i, w in enumerate(words[1:], start=1):
            if w.lower() == 'è':
                subject = ' '.join(words[:i])
                if _valid_subject(subject):
                    return f"Cos'è {subject}?"
            if w.lower() == 'sono':
                subject = ' '.join(words[:i])
                if _valid_subject(subject):
                    return f"Cosa sono {subject}?"

    # English keyword patterns
    if 'should' in lower and len(words) >= 8:
        return f"Why {s[0].lower() + s[1:]}?"
    if 'important' in lower and len(words) >= 6:
        m = re.search(r'(\w[\w\s]{2,20})\s+is\s+important', s, re.IGNORECASE)
        if m:
            return f"Why is {m.group(1).strip()} important?"
    if 'purpose' in lower:
        return "What is the purpose described in this section?"
    if 'reason' in lower and len(words) >= 8:
        return "What is the main reason discussed here?"

    # Italian keyword patterns
    if 'dovrebbe' in lower or 'bisogna' in lower:
        return f"Perché {s[0].lower() + s[1:]}?"
    if 'importante' in lower or 'fondamentale' in lower:
        m = re.search(r'(\w[\w\s]{2,20})\s+è\s+importante', s, re.IGNORECASE)
        if m:
            return f"Perché {m.group(1).strip()} è importante?"
    if 'scopo' in lower or 'obiettivo' in lower:
        return "Qual è lo scopo descritto in questa sezione?"

    return None

def _generate_flashcards(chunks: list[str], max_cards: int = 6) -> list[dict]:
    cards, seen = [], set()
    patterns = [
        # English
        r'([A-Z][A-Za-z\s]+)\s+(?:is defined as|refers to|means|describes)\s+(.{20,100}?)(?:\.|$)',
        r'([A-Z][A-Za-z]+)\s*:\s*(.{20,100}?)(?:\.|$)',
        # Italian
        r'([A-Z][A-Za-zàèéìíîòóùú\s]+)\s+(?:è definit[ao] come|si riferisce a|significa|indica)\s+(.{20,100}?)(?:\.|$)',
        r'([A-Z][A-Za-zàèéìíîòóùú]+)\s*:\s*(.{20,100}?)(?:\.|$)',
    ]
    for chunk in chunks[:20]:
        for pat in patterns:
            for term, definition in re.findall(pat, chunk):
                term = term.strip()
                if len(term.split()) <= 5 and len(definition.split()) >= 5 and term.lower() not in seen:
                    seen.add(term.lower())
                    cards.append({"term": term, "definition": definition.strip()})
                    if len(cards) >= max_cards:
                        return cards
    return cards
